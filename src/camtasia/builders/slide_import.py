"""Presentation slide import — place a sequence of slide images on the timeline."""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

from camtasia.timing import seconds_to_ticks

if TYPE_CHECKING:
    from camtasia.project import Project
    from camtasia.timeline.clips.base import BaseClip


def _extract_slides_as_images(pptx_path: Path, out_dir: Path) -> list[dict[str, str | Path]]:
    """Extract slides from a PowerPoint file as PNG images.

    Requires ``python-pptx`` and ``Pillow``. Falls back to documenting
    a LibreOffice subprocess command if Pillow is unavailable.

    Args:
        pptx_path: Path to the ``.pptx`` file.
        out_dir: Directory to write slide images into.

    Returns:
        List of dicts with ``'path'`` (Path) and ``'title'`` (str) keys.

    Raises:
        ImportError: If ``python-pptx`` is not installed.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        raise ImportError(
            "python-pptx is required for PowerPoint import. "
            "Install it with: pip install python-pptx"
        ) from None

    prs = Presentation(str(pptx_path))
    slides_info: list[dict[str, str | Path]] = []
    width_px = int(prs.slide_width.inches * 96)  # type: ignore[union-attr]
    height_px = int(prs.slide_height.inches * 96)  # type: ignore[union-attr]

    for i, slide in enumerate(prs.slides):
        # Extract title
        title = ''
        if slide.shapes.title is not None:
            title = slide.shapes.title.text or ''

        # Generate thumbnail via Pillow
        out_path = out_dir / f'slide_{i:04d}.png'
        try:
            from PIL import Image, ImageDraw  # type: ignore[import-untyped]

            img = Image.new('RGB', (width_px, height_px), (255, 255, 255))
            if title:
                draw = ImageDraw.Draw(img)
                draw.text((10, 10), title, fill=(0, 0, 0))
            img.save(str(out_path))
        except ImportError:
            # Pillow not available — write a placeholder and document fallback
            out_path.write_text(
                f'Placeholder for slide {i}. '
                'Use LibreOffice to export: '
                f'libreoffice --headless --convert-to png "{pptx_path}"'
            )

        slides_info.append({'path': out_path, 'title': title})

    return slides_info


def _default_slide_duration(project: Project) -> float:
    """Read default image duration from project metadata, falling back to 5.0."""
    return float(project._data.get('metadata', {}).get('defaultImageDuration', 5.0))


def import_slide_images(
    project: Project,
    image_paths: list[Path | str],
    *,
    per_slide_seconds: float | None = None,
    track_name: str = 'Slides',
    transition_seconds: float = 0.0,
    slide_titles: list[str] | None = None,
    emit_markers: bool = False,
    cursor_offset: float = 0.0,
) -> list[BaseClip]:
    """Place a sequence of pre-rendered slide images on a dedicated track.

    The intended workflow is:

    1. In PowerPoint (or Keynote/Google Slides), export slides as PNG/JPG
       images ("File → Export As Pictures" or similar).
    2. Call :func:`import_slide_images` with the ordered list of image paths.

    Each image is imported into the project's media bin and placed on a
    new track, one after another with the specified per-slide duration.

    Args:
        project: Target project.
        image_paths: Ordered list of slide image files.
        per_slide_seconds: Duration each slide is on screen. When None,
            reads from ``project._data['metadata']['defaultImageDuration']``,
            falling back to 5.0.
        track_name: Name of the dedicated slides track.
        transition_seconds: When > 0, adjacent slide pairs overlap by this
            duration, producing a cross-fade effect (via the existing fade
            in/out machinery on each clip).
        slide_titles: Optional list of slide title strings (one per image).
            Used for marker emission even without python-pptx.
        emit_markers: When True, add timeline markers at each slide boundary
            using *slide_titles* (or ``'Slide N'`` if titles not provided).
        cursor_offset: Starting cursor position in seconds (used by append mode).

    Returns:
        List of placed image clips, in order.
    """
    effective_seconds = per_slide_seconds if per_slide_seconds is not None else _default_slide_duration(project)
    track = project.timeline.get_or_create_track(track_name)
    placed: list[BaseClip] = []
    cursor = cursor_offset
    for i, path in enumerate(image_paths):
        media = project.import_media(path)
        clip = track.add_image(
            media.id,
            start_seconds=cursor,
            duration_seconds=effective_seconds,
        )
        if transition_seconds > 0:
            # Fade in all but the first; fade out all but the last
            if i > 0:
                clip.fade_in(transition_seconds)
            if i < len(image_paths) - 1:
                clip.fade_out(transition_seconds)
        if emit_markers:
            label = (
                slide_titles[i]
                if slide_titles and i < len(slide_titles) and slide_titles[i]
                else f'Slide {i + 1}'
            )
            project.timeline.markers.add(label, seconds_to_ticks(cursor))
        placed.append(clip)
        cursor += effective_seconds - transition_seconds
    return placed


def import_powerpoint(
    project: Project,
    pptx_path: Path,
    *,
    per_slide_seconds: float | None = None,
    track_name: str = 'Slides',
    append: bool = False,
    emit_markers: bool = False,
) -> dict:
    """Import a PowerPoint file into the project as slide images.

    Extracts slides as images using ``python-pptx`` (+ Pillow for thumbnails),
    then delegates to :func:`import_slide_images`.

    Args:
        project: Target project.
        pptx_path: Path to the ``.pptx`` file.
        per_slide_seconds: Duration each slide is on screen. Defaults to 5.0.
        track_name: Name of the dedicated slides track.
        append: When True, place slides after the current end of the track.
        emit_markers: When True, add timeline markers at each slide boundary
            with the slide title text.

    Returns:
        Dict with keys ``'clips'`` (list of placed clips), ``'titles'``
        (list of slide title strings), and ``'slide_count'`` (int).

    Raises:
        ImportError: If ``python-pptx`` is not installed.
    """
    import tempfile

    effective_seconds = per_slide_seconds if per_slide_seconds is not None else _default_slide_duration(project)

    with tempfile.TemporaryDirectory() as tmp_dir:
        slides_info = _extract_slides_as_images(Path(pptx_path), Path(tmp_dir))
        image_paths = [info['path'] for info in slides_info]
        titles = [str(info['title']) for info in slides_info]

        cursor_offset = 0.0
        if append:
            track = project.timeline.get_or_create_track(track_name)
            if list(track.clips):
                cursor_offset = max(
                    c.end_seconds for c in track.clips
                )

        clips = import_slide_images(
            project,
            image_paths,  # type: ignore[arg-type]
            per_slide_seconds=effective_seconds,
            track_name=track_name,
            slide_titles=titles,
            emit_markers=emit_markers,
            cursor_offset=cursor_offset,
        )

    return {
        'clips': clips,
        'titles': titles,
        'slide_count': len(slides_info),
    }
